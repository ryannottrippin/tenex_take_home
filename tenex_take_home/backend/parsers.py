import asyncio
import io

import httpx
import openpyxl
from docx import Document
from pypdf import PdfReader
from pptx import Presentation

MAX_CHARS = 10000


async def get_file_content(
    client: httpx.AsyncClient, file: dict, access_token: str
) -> tuple[str | None, list[dict]]:
    """Return (flat_text, sections) for a file.

    sections = [{text, page_label}] where page_label is e.g. "p. 2", "Slide 3", or None.
    flat_text is used for the fallback context-stuffing path.
    sections is used by the vector store for page-aware chunking.
    """
    mime = file.get("mimeType", "")
    file_id = file["id"]
    headers = {"Authorization": f"Bearer {access_token}"}

    try:
        if mime == "application/vnd.google-apps.document":
            response = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}/export",
                params={"mimeType": "text/plain"},
                headers=headers,
            )
            text = response.text[:MAX_CHARS]
            return (text or None), [{"text": text, "page_label": None}]

        elif mime == "application/vnd.google-apps.spreadsheet":
            response = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}/export",
                params={"mimeType": "text/csv"},
                headers=headers,
            )
            text = response.text[:MAX_CHARS]
            return (text or None), [{"text": text, "page_label": None}]

        elif mime == "application/vnd.google-apps.presentation":
            response = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}/export",
                params={"mimeType": "text/plain"},
                headers=headers,
            )
            text = response.text[:MAX_CHARS]
            return (text or None), [{"text": text, "page_label": None}]

        elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            response = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}",
                params={"alt": "media"},
                headers=headers,
            )
            doc = Document(io.BytesIO(response.content))
            text = "\n".join(p.text for p in doc.paragraphs if p.text)
            text = text[:MAX_CHARS]
            return (text or None), [{"text": text, "page_label": None}]

        elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            response = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}",
                params={"alt": "media"},
                headers=headers,
            )
            prs = Presentation(io.BytesIO(response.content))
            sections = []
            for i, slide in enumerate(prs.slides):
                parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                parts.append(para.text.strip())
                    elif shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                parts.append(" | ".join(cells))
                if parts:
                    sections.append({"text": "\n".join(parts), "page_label": f"Slide {i + 1}"})
            flat_text = "\n".join(s["text"] for s in sections)[:MAX_CHARS]
            return (flat_text or None), sections

        elif mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            response = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}",
                params={"alt": "media"},
                headers=headers,
            )
            wb = openpyxl.load_workbook(io.BytesIO(response.content), data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(cell) for cell in row if cell is not None]
                    if cells:
                        parts.append(" | ".join(cells))
            text = "\n".join(parts)[:MAX_CHARS] if parts else None
            return text, ([{"text": text, "page_label": None}] if text else [])

        elif mime == "application/pdf":
            response = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}",
                params={"alt": "media"},
                headers=headers,
            )
            reader = PdfReader(io.BytesIO(response.content))
            sections = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    sections.append({"text": text[:MAX_CHARS], "page_label": f"p. {i + 1}"})
            flat_text = "\n".join(s["text"] for s in sections)[:MAX_CHARS]
            return (flat_text or None), sections

        elif mime in ("text/plain", "text/markdown", "text/csv", "application/json"):
            response = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}",
                params={"alt": "media"},
                headers=headers,
            )
            text = response.text[:MAX_CHARS]
            return (text or None), [{"text": text, "page_label": None}]

        else:
            return None, []

    except Exception:
        return None, []


async def fetch_all_contents(
    client: httpx.AsyncClient, files: list[dict], access_token: str
) -> list[dict]:
    """Fetch content for all files concurrently and return non-empty results."""
    results = await asyncio.gather(
        *[get_file_content(client, f, access_token) for f in files]
    )
    return [
        {"name": f["name"], "id": f["id"], "content": flat_text, "sections": sections}
        for f, (flat_text, sections) in zip(files, results)
        if flat_text
    ]
